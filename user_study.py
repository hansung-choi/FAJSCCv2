from data_maker import *
from loss_maker import *
from optimizer_maker import *
from train import *
from model.model_maker import *
from model_eval import *
import random
import os
import itertools
import random

import os
import csv
import random
import numpy as np
import torch
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ---------------------------
# 1) image save utilities
# ---------------------------
def tensor_to_uint8_img(x_chw: torch.Tensor) -> np.ndarray:
    """
    x_chw: (C,H,W), assumed normalized to [-1,1] like your code
    returns uint8 HxWxC
    """
    x = x_chw.detach().cpu().float()
    x = (x + 1.0) / 2.0
    x = torch.clamp(x, 0.0, 1.0)
    x = (x * 255.0).round().to(torch.uint8)
    x = x.permute(1, 2, 0).contiguous().numpy()
    return x

def save_uint8_png(img_hwc_uint8: np.ndarray, path: str):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    Image.fromarray(img_hwc_uint8).save(path)

# ---------------------------
# 2) generate ref + recon pngs
# ---------------------------
@torch.no_grad()
def save_ref_and_recons(cfg, logger, testloader, model_name_list, num_images: int,
                       out_root: str = "../../test_results/UserStudy",
                       device=None):
    """
    - saves ref images: images/ref/Kodak_Index###.png
    - saves recon images per model: images/recon/<model>/Kodak_Index###.png
    """
    if device is None:
        device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')

    img_root = os.path.join(out_root, "images")
    ref_dir = os.path.join(img_root, "ref")
    recon_root = os.path.join(img_root, "recon")

    os.makedirs(ref_dir, exist_ok=True)
    os.makedirs(recon_root, exist_ok=True)

    # Pre-load all models (one by one) and run inference for each batch index
    # We assume Kodak loader yields deterministic order.
    # We'll cache the refs in memory for num_images only to reuse across models.
    refs = []  # list of torch.Tensor (C,H,W) on CPU
    batches = []

    count = 0
    for images, labels in testloader:
        # images: (B,C,H,W)
        # take first image of batch (your code uses [0] anyway)
        img0 = images[0].clone().cpu()
        refs.append(img0)
        batches.append(images[0:1].clone())  # keep as (1,C,H,W) CPU
        count += 1
        if count >= num_images:
            break

    # Save ref images
    for i, ref_chw in enumerate(refs, start=1):
        stem = f"Kodak_Index{str(i).zfill(3)}"
        ref_path = os.path.join(ref_dir, f"{stem}.png")
        save_uint8_png(tensor_to_uint8_img(ref_chw), ref_path)

    # For each model, load weights and save reconstructions
    for model_name in model_name_list:
        logger.info(f"[UserStudy] Generating reconstructions for model: {model_name}")
        # clone cfg but simplest: temporarily overwrite
        old_name = cfg.model_name
        cfg.model_name = model_name

        # build + load model (reuse your existing save_name logic if needed)
        model = ModelMaker(cfg)
        criterion = LossMaker(cfg)   
        # ---- Load checkpoint like your main() ----
        task = cfg.task_name
        data = cfg.data_info.data_name
        chan_type = cfg.chan_type
        SNR = str(cfg.SNR_info).zfill(3)
        rcpp = str(cfg.rcpp).zfill(3)
        metric = cfg.performance_metric

        save_dir = "../../saved_models/"
        save_name = f"{task}_{data}_{chan_type}_SNR{SNR}_rcpp{rcpp}_{metric}_{model_name}.pt"
        save_name_backup = f"{task}_{data}_{chan_type}_SNR{SNR}_rcpp{rcpp}_{metric}_{model_name}_backup.pt"
        if model_name in ["smallFAJSCCwSA", "baseFAJSCCwSA"]:
            save_name = f"{task}_{data}_{chan_type}_rcpp{rcpp}_{metric}_{model_name}.pt"
            save_name_backup = f"{task}_{data}_{chan_type}_rcpp{rcpp}_{metric}_{model_name}_backup.pt"

        model_path = os.path.join(save_dir, save_name)
        backup_path = os.path.join(save_dir, save_name_backup)
        if not os.path.exists(model_path):
            raise FileNotFoundError(f"Missing checkpoint: {model_path}")

        try:
            model.load_state_dict(torch.load(model_path, map_location="cpu"))
        except Exception:
            model.load_state_dict(torch.load(backup_path, map_location="cpu"))

        model.to(device)
        model.eval()

        model_recon_dir = os.path.join(recon_root, model_name)
        os.makedirs(model_recon_dir, exist_ok=True)

        for i, x_cpu in enumerate(batches, start=1):
            x = x_cpu.to(device)
            y = model(x, SNR_info=cfg.SNR_info)  # (1,C,H,W)
            y0 = y[0].clone().cpu()

            stem = f"Kodak_Index{str(i).zfill(3)}"
            out_path = os.path.join(model_recon_dir, f"{stem}.png")
            save_uint8_png(tensor_to_uint8_img(y0), out_path)

        cfg.model_name = old_name

    logger.info(f"[UserStudy] Saved ref + recon PNGs under: {img_root}")
    return img_root

# ---------------------------
# 3) create trials (pairs) + csv
# ---------------------------
def make_trials(model_name_list, num_images: int, trials_per_image: int = 1, seed: int = 1234):
    """
    Returns list of trials:
      dict(trial_id, stem, model_a, model_b, left_model, right_model, swapped)
    - For each image index, sample model pairs.
    """
    rng = random.Random(seed)
    trials = []
    trial_id = 1

    if len(model_name_list) < 2:
        raise ValueError("Need at least 2 models for pairwise comparison.")

    for i in range(1, num_images + 1):
        stem = f"Kodak_Index{str(i).zfill(3)}"
        for _ in range(trials_per_image):
            a, b = rng.sample(model_name_list, 2)

            # randomize left/right to avoid side bias
            swapped = (rng.random() < 0.5)
            left_model, right_model = (b, a) if swapped else (a, b)

            trials.append({
                "trial_id": trial_id,
                "stem": stem,
                "model_a": a,
                "model_b": b,
                "left_model": left_model,
                "right_model": right_model,
                "swapped": int(swapped),
            })
            trial_id += 1
    return trials

def save_trials_csv(trials, img_root: str, cfg, out_csv_path: str):
    """
    CSV includes separate columns for left/right model names and file paths.
    """
    os.makedirs(os.path.dirname(out_csv_path), exist_ok=True)

    ref_dir = os.path.join(img_root, "ref")
    recon_root = os.path.join(img_root, "recon")

    fieldnames = [
        "trial_id", "stem",
        "task", "data", "chan_type", "SNR_info", "rcpp", "metric",
        "left_model", "right_model", "swapped",
        "ref_path", "left_path", "right_path"
    ]

    with open(out_csv_path, "w", newline="", encoding="utf-8") as f:
        wr = csv.DictWriter(f, fieldnames=fieldnames)
        wr.writeheader()

        for t in trials:
            stem = t["stem"]
            ref_path = os.path.join(ref_dir, f"{stem}.png")
            left_path = os.path.join(recon_root, t["left_model"], f"{stem}.png")
            right_path = os.path.join(recon_root, t["right_model"], f"{stem}.png")

            row = dict(t)
            row.update({
                "task": cfg.task_name,
                "data": cfg.data_info.data_name,
                "chan_type": cfg.chan_type,
                "SNR_info": cfg.SNR_info,
                "rcpp": cfg.rcpp,
                "metric": cfg.performance_metric,
                "ref_path": ref_path,
                "left_path": left_path,
                "right_path": right_path,
            })
            ##drop internal model_a/model_b if you don't want them
            row.pop("model_a", None)
            row.pop("model_b", None)
            wr.writerow(row)

# ---------------------------
# 4) build PPTX
# ---------------------------
def _add_text(slide, text, x, y, w, h, font_size=18, bold=True, align="left"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    else:
        p.alignment = PP_ALIGN.LEFT

def _add_image_centered(slide, img_path, box_x, box_y, box_w, box_h):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size

    box_ar = box_w / box_h
    img_ar = iw / ih if ih else 1.0

    if img_ar >= box_ar:
        w = box_w
        h = box_w / img_ar
    else:
        h = box_h
        w = box_h * img_ar

    x = box_x + (box_w - w) / 2
    y = box_y + (box_h - h) / 2
    slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))

def build_user_study_ppt(trials, img_root: str, out_pptx_path: str):
    import os
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    os.makedirs(os.path.dirname(out_pptx_path), exist_ok=True)

    ref_dir = os.path.join(img_root, "ref")
    recon_root = os.path.join(img_root, "recon")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide_w = 13.333
    slide_h = 7.5


    m = 0.15          # outer margin 
    gap_x = 0.15      # left-right gap 
    gap_y = 0.12      # vertical gap between rows 

    header_y = 0.05
    header_h = 0.45   

    # ---- width constraint from bottom two images ----
    img_w = (slide_w - 2*m - gap_x) / 2

    # ---- height constraint: header + gap + two rows of equal-height images + margins ----
    usable_h = slide_h - (header_y + header_h) - m - gap_y - gap_y  # top margin already in header_y
    #: [header] + gap_y + [ref] + gap_y + [bottom row] + bottom margin(m)
    img_h = (usable_h) / 2

    # positions
    ref_x = (slide_w - img_w) / 2
    ref_y = header_y + header_h + gap_y

    left_x = m
    right_x = m + img_w + gap_x
    bottom_y = ref_y + img_h + gap_y

    for t in trials:
        slide = prs.slides.add_slide(blank)

        trial_id = t["trial_id"]
        stem = t["stem"]
        left_model = t["left_model"]
        right_model = t["right_model"]

        ref_path = os.path.join(ref_dir, f"{stem}.png")
        left_path = os.path.join(recon_root, left_model, f"{stem}.png")
        right_path = os.path.join(recon_root, right_model, f"{stem}.png")

        # Header
        tb = slide.shapes.add_textbox(Inches(m), Inches(header_y), Inches(slide_w - 2*m), Inches(header_h))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"Trial {trial_id:03d} ? Which reconstruction is closer to the reference? (Left / Right)"
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Reference (same size as left/right)
        slide.shapes.add_picture(ref_path, Inches(ref_x), Inches(ref_y), width=Inches(img_w), height=Inches(img_h))

        # Left / Right (same size)
        slide.shapes.add_picture(left_path, Inches(left_x), Inches(bottom_y), width=Inches(img_w), height=Inches(img_h))
        slide.shapes.add_picture(right_path, Inches(right_x), Inches(bottom_y), width=Inches(img_w), height=Inches(img_h))

    prs.save(out_pptx_path)
    
    
def build_user_study_ppt__(trials, img_root: str, out_pptx_path: str): #Previous version
    os.makedirs(os.path.dirname(out_pptx_path), exist_ok=True)

    ref_dir = os.path.join(img_root, "ref")
    recon_root = os.path.join(img_root, "recon")

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Layout (in inches)
    margin = 0.45
    header_y = 0.15
    header_h = 0.65

    gap = 0.25
    ref_label_h = 0.30
    ref_y = header_y + header_h + 0.10
    ref_h = 3.35

    label_y = ref_y + ref_h + gap
    opt_label_h = 0.35

    slide_w = 13.333
    slide_h = 7.5
    half_w = (slide_w - 2 * margin - gap) / 2
    left_x = margin
    right_x = margin + half_w + gap

    opts_y = label_y + opt_label_h + 0.05
    opts_h = slide_h - opts_y - margin

    for t in trials:
        slide = prs.slides.add_slide(blank)

        trial_id = t["trial_id"]
        stem = t["stem"]
        left_model = t["left_model"]
        right_model = t["right_model"]

        ref_path = os.path.join(ref_dir, f"{stem}.png")
        left_path = os.path.join(recon_root, left_model, f"{stem}.png")
        right_path = os.path.join(recon_root, right_model, f"{stem}.png")

        # Header
        _add_text(
            slide,
            f"Trial {trial_id:03d}  ?  Question: Which reconstruction is closer to the reference image? (Left / Right)",
            x=margin, y=header_y, w=slide_w - 2*margin, h=header_h,
            font_size=18, bold=True, align="left"
        )

        # Reference label + image
        _add_text(slide, "Reference", x=margin, y=ref_y - ref_label_h, w=slide_w - 2*margin, h=ref_label_h,
                  font_size=16, bold=True, align="center")
        _add_image_centered(slide, ref_path, box_x=margin, box_y=ref_y, box_w=slide_w - 2*margin, box_h=ref_h)

        # Option labels
        _add_text(slide, "LEFT", x=left_x, y=label_y, w=half_w, h=opt_label_h, font_size=16, bold=True, align="center")
        _add_text(slide, "RIGHT", x=right_x, y=label_y, w=half_w, h=opt_label_h, font_size=16, bold=True, align="center")

        # Option images
        _add_image_centered(slide, left_path, box_x=left_x, box_y=opts_y, box_w=half_w, box_h=opts_h)
        _add_image_centered(slide, right_path, box_x=right_x, box_y=opts_y, box_w=half_w, box_h=opts_h)

    prs.save(out_pptx_path)

# ---------------------------
# 5) one-shot runner
# ---------------------------
def make_user_study_assets(cfg, logger, model_name_list,
                           num_images: int = 60,
                           trials_per_image: int = 1,
                           seed: int = 1234,
                           out_root: str = "../../test_results/UserStudy"):
    """
    1) save ref + recon pngs
    2) create trials
    3) save csv
    4) create ppt
    """
    # build testloader (same way as your main)
    cfg.test_data = "Kodak"
    data_info = DataMaker(cfg)
    testloader = data_info.testloader

    img_root = save_ref_and_recons(cfg, logger, testloader, model_name_list, num_images=num_images, out_root=out_root)
    #trials = make_trials(model_name_list, num_images=num_images, trials_per_image=trials_per_image, seed=seed)
    trials = make_trials_all_pairs(
    model_name_list,
    num_images=num_images,
    seed=seed,
    shuffle_trial_order=True,  # trial shuffle
    shuffle_sides=True         # left, right random trial
    )

    out_csv = os.path.join(out_root, "csv", "user_study_map.csv")
    save_trials_csv(trials, img_root, cfg, out_csv)

    out_pptx = os.path.join(out_root, "ppt", "user_study.pptx")
    build_user_study_ppt(trials, img_root, out_pptx)

    logger.info(f"[UserStudy] PPT saved: {out_pptx}")
    logger.info(f"[UserStudy] CSV saved: {out_csv}")
    return out_pptx, out_csv
    
def make_trials_all_pairs(model_name_list, num_images: int,
                          seed: int = 1234,
                          shuffle_trial_order: bool = True,
                          shuffle_sides: bool = True):
    """
    For each image, generate all unordered model pairs (i<j).
    Then assign them to LEFT/RIGHT (optionally randomized).
    Returns list of trials with explicit left/right model names.
    """
    rng = random.Random(seed)
    pairs = list(itertools.combinations(model_name_list, 2))  # all pairs

    trials = []
    trial_id = 1

    for i in range(1, num_images + 1):
        stem = f"Kodak_Index{str(i).zfill(3)}"

        # Optional: shuffle pair order per image to avoid systematic order effects
        pairs_i = pairs.copy()
        if shuffle_trial_order:
            rng.shuffle(pairs_i)

        for (a, b) in pairs_i:
            if shuffle_sides and (rng.random() < 0.5):
                left_model, right_model = b, a
                swapped = 1
            else:
                left_model, right_model = a, b
                swapped = 0

            trials.append({
                "trial_id": trial_id,
                "stem": stem,
                "left_model": left_model,
                "right_model": right_model,
                "swapped": swapped,
            })
            trial_id += 1

    # Optional: shuffle global trial order across images as well (often good)
    #
    if shuffle_trial_order:
        rng.shuffle(trials)
        for k, t in enumerate(trials, start=1):
            t["trial_id"] = k

    return trials    
    
@hydra.main(version_base = '1.1',config_path="configs",config_name='model_eval')
def main(cfg: DictConfig):
    logger = logging.getLogger(__name__)
    
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    logger.info(f'---------------------------------------------------------------')
    logger.info(f'device: {device}')
    
    hydra_cfg = hydra.core.hydra_config.HydraConfig.get()

    # set random seed number
    random_seed_num = cfg.random_seed
    torch.manual_seed(random_seed_num)
    np.random.seed(random_seed_num)
    random.seed(random_seed_num)

    model_list = ["SwinJSCC", "FAJSCC", "LICRFJSCC", "ResJSCC"]
    make_user_study_assets(
    cfg, logger,
    model_name_list=model_list,
    num_images=10,
    trials_per_image=1,  
    seed=42,
    out_root="../../test_results/UserStudy"
    )







if __name__ == '__main__':
    main()
    
    